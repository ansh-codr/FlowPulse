#!/usr/bin/env python3
"""
FlowPulse 2.0 – PowerPoint Presentation Generator
Generates a professional dark-themed PPT from the project documentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy background
BG_CARD      = RGBColor(0x1A, 0x25, 0x3C)   # Card/content area
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # Cyan accent
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)   # Green accent
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)   # Amber/yellow accent
ACCENT_PINK  = RGBColor(0xF4, 0x72, 0xB6)   # Pink accent
TEXT_WHITE    = RGBColor(0xF1, 0xF5, 0xF9)   # Primary text
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # Secondary/muted text
BORDER_COLOR = RGBColor(0x33, 0x41, 0x55)   # Border/divider
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout for full control
    blank_layout = prs.slide_layouts[6]

    def add_bg(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        # Set corner radius
        if radius:
            shape.adjustments[0] = radius
        else:
            shape.adjustments[0] = 0.02
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                     font_name="Calibri", line_spacing=1.2):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        if line_spacing != 1.0:
            p.line_spacing = Pt(font_size * line_spacing)
        return txBox

    def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                          color=TEXT_WHITE, font_name="Calibri", line_spacing=1.5,
                          alignment=PP_ALIGN.LEFT):
        """lines = list of (text, color, bold, font_size_override)"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(lines):
            if isinstance(item, str):
                txt, clr, bld, fs = item, color, False, font_size
            else:
                txt = item[0]
                clr = item[1] if len(item) > 1 else color
                bld = item[2] if len(item) > 2 else False
                fs  = item[3] if len(item) > 3 else font_size
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = txt
            p.font.size = Pt(fs)
            p.font.color.rgb = clr
            p.font.bold = bld
            p.font.name = font_name
            p.alignment = alignment
            p.space_after = Pt(2)
            p.line_spacing = Pt(fs * line_spacing)
        return txBox

    def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, thickness=3):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_icon_circle(slide, left, top, size, color, label="", label_size=20):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        if label:
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(label_size)
            p.font.color.rgb = TEXT_WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(0)
            tf.paragraphs[0].space_after = Pt(0)
        return shape

    def slide_number_footer(slide, num, total):
        add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
                     f"{num}/{total}", font_size=10, color=TEXT_MUTED,
                     alignment=PP_ALIGN.RIGHT)

    TOTAL_SLIDES = 15

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — TITLE SLIDE
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    # Decorative top accent bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    # Decorative circles
    add_icon_circle(slide, Inches(10.5), Inches(1.0), Inches(2.5), RGBColor(0x38, 0xBD, 0xF8))
    add_icon_circle(slide, Inches(11.2), Inches(2.8), Inches(1.5), RGBColor(0x4A, 0xDE, 0x80))
    add_icon_circle(slide, Inches(9.8), Inches(3.2), Inches(1.0), RGBColor(0xFB, 0xBF, 0x24))

    # Title
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.2),
                 "FlowPulse 2.0", font_size=54, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(1.2), Inches(3.1), Inches(3), ACCENT_BLUE, 4)
    add_text_box(slide, Inches(1.2), Inches(3.4), Inches(8), Inches(0.8),
                 "Browser Extension + Web Dashboard for Productivity Tracking",
                 font_size=22, color=TEXT_MUTED)
    add_text_box(slide, Inches(1.2), Inches(4.4), Inches(8), Inches(0.5),
                 "Powered by Firebase  |  Serverless Architecture  |  Real-Time Analytics",
                 font_size=14, color=ACCENT_BLUE)

    add_text_box(slide, Inches(1.2), Inches(5.8), Inches(8), Inches(0.5),
                 "Project Documentation & Implementation Plan",
                 font_size=16, color=TEXT_MUTED)

    slide_number_footer(slide, 1, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — TABLE OF CONTENTS
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                 "Table of Contents", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE, 3)

    toc_items = [
        ("01", "Project Overview", ACCENT_BLUE),
        ("02", "Problem & Solution", ACCENT_GREEN),
        ("03", "System Architecture", ACCENT_BLUE),
        ("04", "Technology Stack", ACCENT_AMBER),
        ("05", "System Components", ACCENT_PINK),
        ("06", "Database Design", ACCENT_GREEN),
        ("07", "Functional Requirements", ACCENT_BLUE),
        ("08", "Non-Functional Requirements", ACCENT_AMBER),
        ("09", "Development Phases", ACCENT_PINK),
        ("10", "Implementation Plan", ACCENT_GREEN),
        ("11", "Completion Criteria", ACCENT_BLUE),
        ("12", "Folder Structure", ACCENT_AMBER),
        ("13", "Summary & Next Steps", ACCENT_PINK),
    ]

    col1 = toc_items[:7]
    col2 = toc_items[7:]

    for i, (num, title, color) in enumerate(col1):
        y = Inches(1.7) + Inches(i * 0.65)
        add_icon_circle(slide, Inches(1.0), y, Inches(0.45), color, num, 12)
        add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(4), Inches(0.4),
                     title, font_size=17, color=TEXT_WHITE)

    for i, (num, title, color) in enumerate(col2):
        y = Inches(1.7) + Inches(i * 0.65)
        add_icon_circle(slide, Inches(7.0), y, Inches(0.45), color, num, 12)
        add_text_box(slide, Inches(7.7), y + Inches(0.05), Inches(4), Inches(0.4),
                     title, font_size=17, color=TEXT_WHITE)

    slide_number_footer(slide, 2, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — PROJECT OVERVIEW
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                 "Project Overview", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE, 3)

    cards = [
        ("Project Name", "FlowPulse", ACCENT_BLUE),
        ("Type", "Browser Extension +\nWeb Application", ACCENT_GREEN),
        ("Architecture", "Serverless\n(Firebase-based)", ACCENT_AMBER),
        ("Backend", "Firebase Auth +\nFirestore + Functions", ACCENT_PINK),
    ]

    for i, (label, value, color) in enumerate(cards):
        x = Inches(0.8 + i * 3.05)
        card = add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(2.2), BG_CARD, BORDER_COLOR, 0.05)
        add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(0.06), color)
        add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.4),
                     label, font_size=13, color=TEXT_MUTED, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(1.2),
                     value, font_size=20, color=TEXT_WHITE, bold=True)

    # Core Objective box
    add_shape(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.2), BG_CARD, BORDER_COLOR, 0.03)
    add_text_box(slide, Inches(1.1), Inches(4.6), Inches(3), Inches(0.4),
                 "CORE OBJECTIVE", font_size=13, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.1), Inches(5.1), Inches(11.2), Inches(1.2),
                 "Track daily web activity, analyze usage patterns, and provide productivity "
                 "insights through a centralized dashboard using Firebase backend.",
                 font_size=18, color=TEXT_WHITE, line_spacing=1.4)

    slide_number_footer(slide, 3, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — PROBLEM & SOLUTION
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                 "Problem & Solution", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_GREEN, 3)

    # Problem card
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.5), BG_CARD, BORDER_COLOR, 0.04)
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.06), RGBColor(0xEF, 0x44, 0x44))
    add_icon_circle(slide, Inches(1.2), Inches(2.1), Inches(0.6), RGBColor(0xEF, 0x44, 0x44), "!", 22)
    add_text_box(slide, Inches(2.0), Inches(2.15), Inches(4), Inches(0.5),
                 "PROBLEM", font_size=20, color=RGBColor(0xEF, 0x44, 0x44), bold=True)

    problem_lines = [
        ("Users are unaware of how they spend", TEXT_WHITE),
        ("time online", TEXT_WHITE),
        ("", TEXT_WHITE),
        ("No simple, privacy-focused system exists", TEXT_MUTED),
        ("that visualizes browsing behavior", TEXT_MUTED),
        ("across sessions in real-time", TEXT_MUTED),
    ]
    add_multiline_box(slide, Inches(1.2), Inches(3.0), Inches(4.8), Inches(3),
                      problem_lines, font_size=16, line_spacing=1.4)

    # Solution card
    add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.5), BG_CARD, BORDER_COLOR, 0.04)
    add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.06), ACCENT_GREEN)
    add_icon_circle(slide, Inches(7.3), Inches(2.1), Inches(0.6), ACCENT_GREEN, "✓", 22)
    add_text_box(slide, Inches(8.1), Inches(2.15), Inches(4), Inches(0.5),
                 "SOLUTION", font_size=20, color=ACCENT_GREEN, bold=True)

    solution_lines = [
        ("FlowPulse collects browser activity", TEXT_WHITE),
        ("via Chrome extension", TEXT_WHITE),
        ("", TEXT_WHITE),
        ("→  Stores structured logs in Firebase", ACCENT_BLUE),
        ("→  Processes analytics via Cloud Functions", ACCENT_BLUE),
        ("→  Displays insights in React dashboard", ACCENT_BLUE),
    ]
    add_multiline_box(slide, Inches(7.3), Inches(3.0), Inches(4.8), Inches(3),
                      solution_lines, font_size=16, line_spacing=1.4)

    slide_number_footer(slide, 4, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — SYSTEM ARCHITECTURE
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                 "System Architecture", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE, 3)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                 "Serverless Architecture  –  End-to-End Data Flow",
                 font_size=16, color=TEXT_MUTED)

    # Architecture flow boxes
    arch_items = [
        ("🌐", "Browser\nExtension", "Tab tracking\nIdle detection\nDomain extraction", ACCENT_BLUE),
        ("🔐", "Firebase\nAuth", "Google Sign-In\nUser isolation\nSession mgmt", ACCENT_GREEN),
        ("🗄️", "Cloud\nFirestore", "Activity logs\nUser settings\nDaily stats", ACCENT_AMBER),
        ("⚡", "Cloud\nFunctions", "Daily aggregation\nProductivity score\nPattern detection", ACCENT_PINK),
        ("📊", "Web\nDashboard", "Charts & graphs\nSettings UI\nReal-time sync", ACCENT_BLUE),
    ]

    for i, (icon, title, desc, color) in enumerate(arch_items):
        x = Inches(0.5 + i * 2.55)
        y = Inches(2.3)

        # Card
        add_shape(slide, x, y, Inches(2.3), Inches(4.2), BG_CARD, BORDER_COLOR, 0.05)
        add_shape(slide, x, y, Inches(2.3), Inches(0.06), color)

        # Icon
        add_text_box(slide, x, y + Inches(0.3), Inches(2.3), Inches(0.7),
                     icon, font_size=36, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.15), y + Inches(1.1), Inches(2.0), Inches(0.8),
                     title, font_size=17, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        desc_lines = [(line, TEXT_MUTED) for line in desc.split("\n")]
        add_multiline_box(slide, x + Inches(0.15), y + Inches(2.2), Inches(2.0), Inches(1.8),
                          desc_lines, font_size=12, line_spacing=1.5, alignment=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < len(arch_items) - 1:
            arrow_x = x + Inches(2.35)
            add_text_box(slide, arrow_x, y + Inches(1.5), Inches(0.3), Inches(0.5),
                         "→", font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 5, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — TECHNOLOGY STACK
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                 "Technology Stack", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_AMBER, 3)

    tech_groups = [
        ("Frontend (Web App)", [
            "React", "Firebase SDK", "Recharts", "Tailwind CSS", "Vite"
        ], ACCENT_BLUE),
        ("Browser Extension", [
            "JavaScript", "Chrome APIs (MV3)", "Firebase SDK"
        ], ACCENT_GREEN),
        ("Backend Services", [
            "Firebase Authentication", "Cloud Firestore", "Cloud Functions (Node.js)", "Firebase Hosting"
        ], ACCENT_AMBER),
    ]

    for i, (group_title, techs, color) in enumerate(tech_groups):
        x = Inches(0.8 + i * 4.1)
        add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(5.0), BG_CARD, BORDER_COLOR, 0.04)
        add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(0.06), color)

        add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3.2), Inches(0.5),
                     group_title, font_size=18, color=color, bold=True)

        lines = []
        for tech in techs:
            lines.append(("●  " + tech, TEXT_WHITE))
        add_multiline_box(slide, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(3.5),
                          lines, font_size=16, line_spacing=1.8)

    slide_number_footer(slide, 6, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — SYSTEM COMPONENTS (EXTENSION)
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "System Components — Browser Extension",
                 font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_PINK, 3)

    # Responsibilities
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.0), BG_CARD, BORDER_COLOR, 0.04)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5),
                 "RESPONSIBILITIES", font_size=14, color=ACCENT_PINK, bold=True)
    resp_lines = [
        ("●  Detect active tab and track time per domain", TEXT_WHITE),
        ("●  Extract domain and page title from tabs", TEXT_WHITE),
        ("●  Calculate browsing duration per site", TEXT_WHITE),
        ("●  Send activity logs to Firestore in batches", TEXT_WHITE),
        ("●  Sync user settings (tracking toggle, blocked sites)", TEXT_WHITE),
    ]
    add_multiline_box(slide, Inches(1.1), Inches(2.5), Inches(5.2), Inches(3.5),
                      resp_lines, font_size=15, line_spacing=1.8)

    # Key Modules
    add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(5.0), BG_CARD, BORDER_COLOR, 0.04)
    add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5), Inches(0.5),
                 "KEY MODULES", font_size=14, color=ACCENT_BLUE, bold=True)

    modules = [
        ("Background Script", "Tab tracking, idle detection, alarm-based flush"),
        ("Content Script", "Title extraction, focus/blur detection"),
        ("Popup UI", "Quick stats, sign-in/out, toggle tracking"),
        ("Auth Handler", "Google login flow, session management"),
    ]

    for i, (name, desc) in enumerate(modules):
        y = Inches(2.5) + Inches(i * 1.0)
        add_icon_circle(slide, Inches(7.3), y, Inches(0.4), ACCENT_BLUE, str(i + 1), 13)
        add_text_box(slide, Inches(7.9), y - Inches(0.05), Inches(4.2), Inches(0.35),
                     name, font_size=15, color=TEXT_WHITE, bold=True)
        add_text_box(slide, Inches(7.9), y + Inches(0.3), Inches(4.2), Inches(0.35),
                     desc, font_size=12, color=TEXT_MUTED)

    slide_number_footer(slide, 7, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — SYSTEM COMPONENTS (BACKEND + DASHBOARD)
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "System Components — Backend & Dashboard",
                 font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_GREEN, 3)

    # Firebase Backend
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.0), BG_CARD, BORDER_COLOR, 0.04)
    add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.06), ACCENT_AMBER)
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5),
                 "FIREBASE BACKEND", font_size=14, color=ACCENT_AMBER, bold=True)

    backend_lines = [
        ("Authentication", ACCENT_BLUE, True),
        ("  Google Sign-In, user-based data isolation", TEXT_MUTED),
        ("", TEXT_WHITE),
        ("Cloud Firestore", ACCENT_BLUE, True),
        ("  Activity logs, user settings, daily stats", TEXT_MUTED),
        ("", TEXT_WHITE),
        ("Cloud Functions", ACCENT_BLUE, True),
        ("  Daily aggregation, productivity scoring", TEXT_MUTED),
        ("  Pattern detection, weekly trends", TEXT_MUTED),
    ]
    add_multiline_box(slide, Inches(1.1), Inches(2.5), Inches(5.2), Inches(4),
                      backend_lines, font_size=14, line_spacing=1.5)

    # Web Dashboard
    add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(5.0), BG_CARD, BORDER_COLOR, 0.04)
    add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.06), ACCENT_GREEN)
    add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5), Inches(0.5),
                 "WEB DASHBOARD FEATURES", font_size=14, color=ACCENT_GREEN, bold=True)

    dashboard_lines = [
        ("●  Daily activity chart (time per domain)", TEXT_WHITE),
        ("●  Category breakdown (productive/distraction)", TEXT_WHITE),
        ("●  Weekly trends (line/area charts)", TEXT_WHITE),
        ("●  Productivity score gauge", TEXT_WHITE),
        ("●  Enable/Disable tracking toggle", TEXT_WHITE),
        ("●  Manage blocked sites list", TEXT_WHITE),
        ("●  Real-time data updates", TEXT_WHITE),
    ]
    add_multiline_box(slide, Inches(7.2), Inches(2.5), Inches(5.0), Inches(4),
                      dashboard_lines, font_size=14, line_spacing=1.7)

    slide_number_footer(slide, 8, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — DATABASE DESIGN
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                 "Database Design", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_GREEN, 3)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
                 "Firestore Collections Structure — Nested under users/{userId}",
                 font_size=15, color=TEXT_MUTED)

    # Collection cards
    collections = [
        ("users/{userId}", [
            ("email", "string"),
            ("displayName", "string"),
            ("photoURL", "string"),
            ("createdAt", "timestamp"),
        ], ACCENT_BLUE),
        ("settings", [
            ("trackingEnabled", "boolean"),
            ("blockedDomains", "string[]"),
            ("timezone", "string"),
        ], ACCENT_GREEN),
        ("activityLogs", [
            ("url", "string"),
            ("domain", "string"),
            ("title", "string"),
            ("category", "string"),
            ("startTime", "timestamp"),
            ("endTime", "timestamp"),
            ("duration", "number (sec)"),
        ], ACCENT_AMBER),
        ("dailyStats", [
            ("totalDuration", "number"),
            ("productiveTime", "number"),
            ("distractionTime", "number"),
            ("topDomains", "map[]"),
            ("focusScore", "number 0-100"),
            ("peakHour", "number 0-23"),
        ], ACCENT_PINK),
    ]

    for i, (coll_name, fields, color) in enumerate(collections):
        x = Inches(0.5 + i * 3.15)
        card_h = Inches(4.8)
        add_shape(slide, x, Inches(2.0), Inches(2.95), card_h, BG_CARD, BORDER_COLOR, 0.04)
        add_shape(slide, x, Inches(2.0), Inches(2.95), Inches(0.06), color)

        add_text_box(slide, x + Inches(0.15), Inches(2.2), Inches(2.65), Inches(0.45),
                     coll_name, font_size=13, color=color, bold=True)

        lines = []
        for fname, ftype in fields:
            lines.append((f"  {fname}", TEXT_WHITE, False, 12))
            lines.append((f"     {ftype}", TEXT_MUTED, False, 10))
        add_multiline_box(slide, x + Inches(0.15), Inches(2.7), Inches(2.65), Inches(3.8),
                          lines, font_size=12, line_spacing=1.3)

    slide_number_footer(slide, 9, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — FUNCTIONAL REQUIREMENTS
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Functional Requirements", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE, 3)

    frs = [
        ("FR-01", "User can sign in securely via Google Login"),
        ("FR-02", "Extension detects active tab and tracks browsing duration"),
        ("FR-03", "Activity logs are stored per user in Firestore"),
        ("FR-04", "Dashboard displays real-time daily analytics"),
        ("FR-05", "Dashboard displays weekly trend charts"),
        ("FR-06", "Dashboard shows productivity score"),
        ("FR-07", "User can enable/disable tracking from dashboard and extension"),
        ("FR-08", "User can manage a blocked-sites list"),
        ("FR-09", "System aggregates daily usage via Cloud Functions"),
        ("FR-10", "Settings sync in real-time between extension and dashboard"),
    ]

    add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2), BG_CARD, BORDER_COLOR, 0.03)

    for i, (code, desc) in enumerate(frs):
        y = Inches(1.95) + Inches(i * 0.48)
        color = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_AMBER, ACCENT_PINK, ACCENT_BLUE][i % 5]
        add_shape(slide, Inches(1.1), y, Inches(0.85), Inches(0.35), color)
        add_text_box(slide, Inches(1.12), y + Inches(0.01), Inches(0.85), Inches(0.35),
                     code, font_size=11, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(2.15), y + Inches(0.01), Inches(9.8), Inches(0.35),
                     desc, font_size=15, color=TEXT_WHITE)

    slide_number_footer(slide, 10, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 11 — NON-FUNCTIONAL REQUIREMENTS
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Non-Functional Requirements", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_AMBER, 3)

    nfrs = [
        ("NFR-01", "Secure User Isolation", "Each user sees only their own data", "🔒"),
        ("NFR-02", "Scalable Design", "Firestore sub-collections for efficient queries", "📈"),
        ("NFR-03", "Low Latency Logging", "Batch writes, minimal network calls", "⚡"),
        ("NFR-04", "Real-Time Sync", "Firestore onSnapshot listeners for settings", "🔄"),
        ("NFR-05", "Minimal Impact", "Background script optimized for performance", "🪶"),
        ("NFR-06", "Responsive UI", "Dashboard works on desktop and tablet", "📱"),
        ("NFR-07", "Data Validation", "Firestore security rules enforce schema", "✅"),
    ]

    for i, (code, title, desc, icon) in enumerate(nfrs):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.7 + row * 1.35)

        add_shape(slide, x, y, Inches(5.9), Inches(1.15), BG_CARD, BORDER_COLOR, 0.04)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                     icon, font_size=22, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.75), y + Inches(0.1), Inches(1.0), Inches(0.35),
                     code, font_size=11, color=ACCENT_AMBER, bold=True)
        add_text_box(slide, x + Inches(1.75), y + Inches(0.1), Inches(3.8), Inches(0.35),
                     title, font_size=15, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x + Inches(0.75), y + Inches(0.55), Inches(4.8), Inches(0.45),
                     desc, font_size=12, color=TEXT_MUTED)

    slide_number_footer(slide, 11, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 12 — DEVELOPMENT PHASES (1-3)
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Development Phases (1–3)", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_PINK, 3)

    phases_1 = [
        ("Phase 1", "Project Setup", [
            "Create Firebase project",
            "Enable Auth (Google provider)",
            "Initialize Cloud Firestore",
            "Define security rules",
            "Set up React + Vite + Tailwind",
            "Set up extension scaffold (MV3)",
            "Configure Firebase SDK",
        ], ACCENT_BLUE),
        ("Phase 2", "Extension Development", [
            "Active tab detection (chrome.tabs)",
            "Idle state detection (chrome.idle)",
            "Track duration per tab session",
            "Extract domain + title from URLs",
            "Batch logs → push to Firestore",
            "Google Sign-In in popup",
            "Build popup UI",
            "Sync settings from Firestore",
        ], ACCENT_GREEN),
        ("Phase 3", "Web Dashboard", [
            "Firebase Auth (login/logout)",
            "Dashboard layout + sidebar",
            "Daily activity chart",
            "Category breakdown (PieChart)",
            "Weekly trends (LineChart)",
            "Productivity score gauge",
            "Settings page",
            "Real-time data (onSnapshot)",
        ], ACCENT_AMBER),
    ]

    for i, (phase, title, tasks, color) in enumerate(phases_1):
        x = Inches(0.6 + i * 4.15)
        add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(5.2), BG_CARD, BORDER_COLOR, 0.04)
        add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(0.06), color)

        add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(0.35),
                     phase, font_size=12, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.25), Inches(3.5), Inches(0.4),
                     title, font_size=18, color=TEXT_WHITE, bold=True)

        lines = [("□  " + t, TEXT_MUTED) for t in tasks]
        add_multiline_box(slide, x + Inches(0.2), Inches(2.8), Inches(3.5), Inches(3.8),
                          lines, font_size=12, line_spacing=1.6)

    slide_number_footer(slide, 12, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 13 — DEVELOPMENT PHASES (4-6)
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Development Phases (4–6)", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_PINK, 3)

    phases_2 = [
        ("Phase 4", "Backend Logic", [
            "Scheduled daily aggregation function",
            "Compute productivity score per day",
            "Compute top domains + peak hour",
            "Write results to dailyStats",
            "Pattern detection for weekly trends",
        ], ACCENT_PINK),
        ("Phase 5", "Optimization & Security", [
            "Firestore security rules (user isolation)",
            "Input validation in Cloud Functions",
            "Rate limiting for activity writes",
            "Performance tuning (batch writes)",
            "Error handling and logging",
        ], RGBColor(0xA7, 0x8B, 0xFA)),
        ("Phase 6", "Deployment", [
            "Package Chrome extension",
            "Deploy web app to Firebase Hosting",
            "Deploy Cloud Functions to prod",
            "Configure production env vars",
            "Final end-to-end testing",
        ], ACCENT_BLUE),
    ]

    for i, (phase, title, tasks, color) in enumerate(phases_2):
        x = Inches(0.6 + i * 4.15)
        add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(4.5), BG_CARD, BORDER_COLOR, 0.04)
        add_shape(slide, x, Inches(1.7), Inches(3.9), Inches(0.06), color)

        add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(0.35),
                     phase, font_size=12, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(2.25), Inches(3.5), Inches(0.4),
                     title, font_size=18, color=TEXT_WHITE, bold=True)

        lines = [("□  " + t, TEXT_MUTED) for t in tasks]
        add_multiline_box(slide, x + Inches(0.2), Inches(2.8), Inches(3.5), Inches(3.2),
                          lines, font_size=13, line_spacing=1.7)

    slide_number_footer(slide, 13, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 14 — COMPLETION CRITERIA
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                 "Completion Criteria", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_BLUE, 3)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
                 "Project is considered COMPLETE when all criteria are met:",
                 font_size=15, color=TEXT_MUTED)

    criteria = [
        "Users can sign in via Google Login",
        "Extension tracks browsing activity accurately",
        "Activity logs are stored in Firestore per user",
        "Dashboard displays daily and weekly insights",
        "Productivity score is calculated and displayed",
        "Settings sync in real-time (extension ↔ dashboard)",
        "Data is secure and isolated per user",
        "System performs reliably under multiple users",
        "Extension is packaged and loadable",
        "Web app is deployed and accessible",
    ]

    add_shape(slide, Inches(2.0), Inches(2.0), Inches(9.3), Inches(5.0), BG_CARD, BORDER_COLOR, 0.04)

    for i, criterion in enumerate(criteria):
        y = Inches(2.3) + Inches(i * 0.45)
        color = ACCENT_GREEN
        add_text_box(slide, Inches(2.5), y, Inches(0.4), Inches(0.35),
                     "✓", font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(3.0), y + Inches(0.02), Inches(7.8), Inches(0.35),
                     criterion, font_size=15, color=TEXT_WHITE)

    slide_number_footer(slide, 14, TOTAL_SLIDES)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 15 — THANK YOU
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    add_text_box(slide, Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(1.2),
                 "Thank You", font_size=54, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(5.5), Inches(3.5), Inches(2.3), ACCENT_BLUE, 4)

    add_text_box(slide, Inches(0), Inches(3.8), SLIDE_WIDTH, Inches(0.7),
                 "FlowPulse 2.0", font_size=28, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER, bold=True)
    add_text_box(slide, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.5),
                 "A Smarter Way to Understand Your Browsing Habits",
                 font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Decorative circles
    add_icon_circle(slide, Inches(1.5), Inches(5.5), Inches(1.2), ACCENT_BLUE)
    add_icon_circle(slide, Inches(3.0), Inches(6.0), Inches(0.8), ACCENT_GREEN)
    add_icon_circle(slide, Inches(9.5), Inches(5.5), Inches(1.2), ACCENT_PINK)
    add_icon_circle(slide, Inches(11.0), Inches(6.0), Inches(0.8), ACCENT_AMBER)

    add_text_box(slide, Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.5),
                 "Questions?", font_size=20, color=TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)

    slide_number_footer(slide, 15, TOTAL_SLIDES)

    # ── Save ──────────────────────────────────────────────────────────────────
    output_path = "/Users/havocerebus/Documents/Current_Project_Working/FlowPulse/FlowPulse_2.0_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   {TOTAL_SLIDES} slides generated")


if __name__ == "__main__":
    create_presentation()
